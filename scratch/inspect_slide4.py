import sys
from pptx import Presentation

# Reconfigure stdout to use UTF-8
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation("TalentDNA_Redrob_Submission.pptx")
slide = prs.slides[3] # Slide 4 is index 3

print("Slide 4 Shapes:")
for idx, shape in enumerate(slide.shapes):
    has_text = shape.has_text_frame
    text = shape.text_frame.text if has_text else ""
    print(f"Index: {idx} | Shape Name: '{shape.name}' | Type: {shape.shape_type} | HasText: {has_text}")
    if has_text and text.strip():
        print(f"  Text: {repr(text)}")
