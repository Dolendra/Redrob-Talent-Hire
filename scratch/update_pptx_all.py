import sys
from pptx import Presentation

# Reconfigure stdout to use UTF-8
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

def update_shape_text(shape, new_text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    
    # Save formatting of the first run of the first paragraph if it exists
    font_name = None
    font_size = None
    bold = None
    color = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        font_name = run.font.name
        font_size = run.font.size
        bold = run.font.bold
        try:
            if run.font.color and hasattr(run.font.color, 'rgb'):
                color = run.font.color.rgb
        except Exception:
            pass
        
    tf.text = new_text
    
    # Re-apply formatting to all runs if we saved it
    if font_name or font_size or bold or color:
        for p in tf.paragraphs:
            for r in p.runs:
                if font_name:
                    r.font.name = font_name
                if font_size:
                    r.font.size = font_size
                if bold is not None:
                    r.font.bold = bold
                if color and hasattr(r.font.color, 'rgb'):
                    try:
                        r.font.color.rgb = color
                    except Exception:
                        pass

def main():
    prs = Presentation("TalentDNA_Redrob_Submission.pptx")
    
    # ------------------ Slide 2: Solution Overview ------------------
    slide2 = prs.slides[1]
    # Update Contextual Understanding description
    update_shape_text(slide2.shapes[9], "Our local LLM understands semantic structure, completely bypassing casing, punctuation, and typographical differences (e.g., accurately mapping 'Sentence Transformers' or 'sentence-transformers' to the same unified canonical schema item).")
    print("Slide 2 updated.")
    
    # ------------------ Slide 4: Ranking Methodology ------------------
    slide4 = prs.slides[3]
    # Header subtext
    update_shape_text(slide4.shapes[2], "Five scoring components → weighted composite → O(N log K) min-heap")
    # Weights percentages
    update_shape_text(slide4.shapes[4], "15%")
    update_shape_text(slide4.shapes[9], "30%")
    update_shape_text(slide4.shapes[14], "20%")
    update_shape_text(slide4.shapes[24], "−15%")
    print("Slide 4 updated.")
    
    prs.save("TalentDNA_Redrob_Submission.pptx")
    print("Successfully updated TalentDNA_Redrob_Submission.pptx")

if __name__ == "__main__":
    main()
